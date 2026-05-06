#!/usr/bin/env python3
"""
make_materials.py — generate per-unit placeholder presentation and
worksheet artefacts for the Materials hub, plus matching thumbnails,
plus a frontmatter update on each unit `.md` adding the file paths.

Outputs land under:
  static/materials/presentations/unit<NN>_<slug>.pptx
  static/materials/presentations/unit<NN>_<slug>.png
  static/materials/worksheets/unit<NN>_<slug>.pdf
  static/materials/worksheets/unit<NN>_<slug>.png

Each binary is a real, valid one-slide / one-page document — they
open in their respective applications. The PNG thumbnails are
drawn directly with Pillow (independent of the .pptx/.pdf renderers
to keep this script Windows-friendly with no Poppler dependency).

Run from repo root:  python scripts/make_materials.py
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from PIL import Image, ImageDraw, ImageFont

REPO = Path(__file__).resolve().parent.parent
CONTENT = REPO / "content"
OUT_PRES = REPO / "static" / "materials" / "presentations"
OUT_WORK = REPO / "static" / "materials" / "worksheets"

ACCENT = RGBColor(0x1A, 0x73, 0xE8)
ACCENT_HEX = "#1A73E8"

UNIT_RE = re.compile(r"^unit(\d{2})_([a-z0-9_-]+)\.md$")


# ---------------------------- font helpers -----------------------------------

def _font(size: int) -> ImageFont.ImageFont:
    """Pick a readable system font; fall back to PIL's default."""
    candidates = [
        "C:/Windows/Fonts/segoeui.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/Library/Fonts/Arial.ttf",
    ]
    for c in candidates:
        if Path(c).exists():
            return ImageFont.truetype(c, size)
    return ImageFont.load_default()


# ---------------------------- generators -------------------------------------

def render_thumbnail(path: Path, kind: str, level: str, unit_nr: int, title: str) -> None:
    """Draw a 1280×720 thumbnail with the unit title + CEFR badge."""
    W, H = 1280, 720
    img = Image.new("RGB", (W, H), color=(255, 255, 255))
    d = ImageDraw.Draw(img)

    # Top stripe
    d.rectangle([(0, 0), (W, 100)], fill=(26, 115, 232))
    d.text((40, 28), f"DaF · GER {level}", font=_font(40), fill=(255, 255, 255))
    d.text(
        (W - 360, 30),
        f"Einheit {unit_nr:02d}",
        font=_font(36),
        fill=(255, 255, 255),
    )

    # Title (wrapped to fit)
    lines = wrap_text(title, _font(54), W - 120)
    y = 200
    for ln in lines[:4]:
        d.text((60, y), ln, font=_font(54), fill=(20, 20, 20))
        y += 70

    # Footer label
    label = "Foliensatz" if kind == "presentation" else "Arbeitsblatt"
    d.text(
        (60, H - 110),
        f"{label} — Platzhalter",
        font=_font(32),
        fill=(100, 100, 100),
    )
    d.text(
        (60, H - 70),
        "S. Le Boulanger · CC-BY-SA 4.0",
        font=_font(24),
        fill=(120, 120, 120),
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, "PNG", optimize=True)


def wrap_text(text: str, font: ImageFont.ImageFont, max_w: int) -> list[str]:
    """Greedy word-wrap that respects pixel width."""
    words, lines, cur = text.split(), [], ""
    for w in words:
        trial = (cur + " " + w).strip()
        if font.getlength(trial) <= max_w:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


def make_pptx(path: Path, level: str, unit_nr: int, title: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Top accent bar
    from pptx.shapes.autoshape import Shape  # noqa: F401
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.text_frame.text = f"DaF · GER {level} · Einheit {unit_nr:02d}"
    p = bar.text_frame.paragraphs[0]
    for r in p.runs:
        r.font.size = Pt(24)
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.0), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(44)

    # Placeholder caption
    cap = slide.shapes.add_textbox(Inches(0.6), Inches(5.5), Inches(12.0), Inches(1.5))
    cap.text_frame.text = (
        "Foliensatz — Platzhalter. Echter Foliensatz folgt iterativ.\n"
        "S. Le Boulanger · CC-BY-SA 4.0"
    )
    for para in cap.text_frame.paragraphs:
        for r in para.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # PDF metadata equivalent: presentation core_properties
    prs.core_properties.title = title
    prs.core_properties.author = "S. Le Boulanger"
    prs.core_properties.subject = f"DaF Einheit {unit_nr:02d} ({level})"

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)


def make_pdf(path: Path, level: str, unit_nr: int, title: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    c = canvas.Canvas(str(path), pagesize=A4)
    c.setAuthor("S. Le Boulanger")
    c.setTitle(title)
    c.setSubject(f"DaF Einheit {unit_nr:02d} ({level}) — Arbeitsblatt-Platzhalter")

    W, H = A4
    # Top bar
    c.setFillColorRGB(0x1A / 255, 0x73 / 255, 0xE8 / 255)
    c.rect(0, H - 70, W, 70, stroke=0, fill=1)
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 14)
    c.drawString(40, H - 45, f"DaF · GER {level} · Einheit {unit_nr:02d}")
    c.drawRightString(W - 40, H - 45, "Arbeitsblatt — Platzhalter")

    # Title
    c.setFillColorRGB(0.1, 0.1, 0.1)
    c.setFont("Helvetica-Bold", 22)
    c.drawString(40, H - 130, title[:80])

    # Body
    c.setFont("Helvetica", 12)
    c.setFillColorRGB(0.3, 0.3, 0.3)
    body = [
        "Echtes Arbeitsblatt folgt iterativ.",
        "",
        "Dieses PDF ist ein gültiger einseitiger Platzhalter, der von der",
        "automatischen Materialien-Übersicht eingebunden wird. Sobald das",
        "endgültige Arbeitsblatt vorliegt, wird diese Datei überschrieben.",
        "",
        "© S. Le Boulanger · CC-BY-SA 4.0",
    ]
    y = H - 180
    for line in body:
        c.drawString(40, y, line)
        y -= 18

    c.showPage()
    c.save()


# ---------------------------- frontmatter rewriter ---------------------------

FM_RE = re.compile(r"^---\n(.*?)\n---\n", re.S)


def update_unit_frontmatter(md: Path, level: str, unit_nr: int, slug: str) -> bool:
    text = md.read_text(encoding="utf-8")
    m = FM_RE.match(text)
    if not m:
        print(f"  skip  {md.relative_to(REPO)} (no frontmatter)")
        return False
    fm = yaml.safe_load(m.group(1)) or {}

    base = f"unit{unit_nr:02d}_{slug}"
    fm["presentation"] = {
        "file": f"/materials/presentations/{base}.pptx",
        "thumbnail": f"/materials/presentations/{base}.png",
    }
    fm["worksheet"] = {
        "file": f"/materials/worksheets/{base}.pdf",
        "thumbnail": f"/materials/worksheets/{base}.png",
    }

    new_fm = yaml.safe_dump(fm, allow_unicode=True, sort_keys=False).strip()
    md.write_text(f"---\n{new_fm}\n---\n{text[m.end():]}", encoding="utf-8")
    return True


# ---------------------------- driver -----------------------------------------

def iter_units():
    for course in sorted(CONTENT.glob("kurs_*")):
        units_dir = course / "units"
        if not units_dir.is_dir():
            continue
        for md in sorted(units_dir.glob("unit*.md")):
            m = UNIT_RE.match(md.name)
            if not m:
                continue
            unit_nr = int(m.group(1))
            slug = m.group(2)
            level = course.name.removeprefix("kurs_").upper()
            yield md, level, unit_nr, slug


def main() -> int:
    n = 0
    for md, level, unit_nr, slug in iter_units():
        text = md.read_text(encoding="utf-8")
        m = FM_RE.match(text)
        fm = yaml.safe_load(m.group(1)) if m else {}
        title = fm.get("title") or md.stem

        base = f"unit{unit_nr:02d}_{slug}"
        pptx_path = OUT_PRES / f"{base}.pptx"
        pres_thumb = OUT_PRES / f"{base}.png"
        pdf_path = OUT_WORK / f"{base}.pdf"
        work_thumb = OUT_WORK / f"{base}.png"

        make_pptx(pptx_path, level, unit_nr, title)
        render_thumbnail(pres_thumb, "presentation", level, unit_nr, title)
        make_pdf(pdf_path, level, unit_nr, title)
        render_thumbnail(work_thumb, "worksheet", level, unit_nr, title)
        update_unit_frontmatter(md, level, unit_nr, slug)
        print(f"  ok  {level} {unit_nr:02d}  {slug}")
        n += 1

    print(f"\n{n} units processed.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
